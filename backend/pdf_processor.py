import os
import pypdf
import pdfplumber
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
import io
import shutil
from typing import List, Dict, Any

class PDFProcessor:
    @staticmethod
    def get_metadata(pdf_path: str) -> Dict[str, Any]:
        """Extract metadata from the PDF using pypdf."""
        reader = pypdf.PdfReader(pdf_path)
        metadata = reader.metadata
        page_count = len(reader.pages)
        is_encrypted = reader.is_encrypted
        
        title = ""
        author = ""
        subject = ""
        keywords = ""
        creator = ""
        producer = ""
        creation_date = ""
        mod_date = ""
        
        if metadata:
            title = metadata.get("/Title", "")
            author = metadata.get("/Author", "")
            subject = metadata.get("/Subject", "")
            keywords = metadata.get("/Keywords", "")
            creator = metadata.get("/Creator", "")
            producer = metadata.get("/Producer", "")
            creation_date = metadata.get("/CreationDate", "")
            mod_date = metadata.get("/ModDate", "")
            
        return {
            "title": str(title),
            "author": str(author),
            "subject": str(subject),
            "keywords": str(keywords),
            "creator": str(creator),
            "producer": str(producer),
            "creationDate": str(creation_date),
            "modDate": str(mod_date),
            "pages": page_count,
            "isEncrypted": is_encrypted
        }

    @staticmethod
    def extract_text_by_page(pdf_path: str) -> List[Dict[str, Any]]:
        """Extract plain text page-by-page using pypdf."""
        reader = pypdf.PdfReader(pdf_path)
        pages_text = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            pages_text.append({
                "page": i + 1,
                "text": text
            })
        return pages_text

    @staticmethod
    def extract_detailed_text(pdf_path: str) -> List[Dict[str, Any]]:
        """Extract text with coordinates, fonts, and sizes using pdfplumber."""
        detailed_data = []
        with pdfplumber.open(pdf_path) as pdf:
            for page_idx, page in enumerate(pdf.pages):
                page_width = float(page.width)
                page_height = float(page.height)
                
                words = page.extract_words()
                page_blocks = []
                for w in words:
                    page_blocks.append({
                        "text": w["text"],
                        "bbox": [float(w["x0"]), float(w["top"]), float(w["x1"]), float(w["bottom"])],
                        "font": w.get("fontname", "Helvetica"),
                        "size": float(w.get("size", 10.0)),
                        "color": 0 # Default black
                    })
                
                detailed_data.append({
                    "page": page_idx + 1,
                    "width": page_width,
                    "height": page_height,
                    "elements": page_blocks
                })
        return detailed_data

    @staticmethod
    def extract_tables(pdf_path: str) -> List[Dict[str, Any]]:
        """Extract tables from the PDF using pdfplumber."""
        tables_data = []
        with pdfplumber.open(pdf_path) as pdf:
            for page_idx, page in enumerate(pdf.pages):
                tables = page.extract_tables()
                for table_idx, table in enumerate(tables):
                    tables_data.append({
                        "page": page_idx + 1,
                        "table_index": table_idx + 1,
                        "data": table
                    })
        return tables_data

    @staticmethod
    def extract_images(pdf_path: str, output_dir: str) -> List[Dict[str, Any]]:
        """Extract images from the PDF using pypdf."""
        reader = pypdf.PdfReader(pdf_path)
        image_info = []
        os.makedirs(output_dir, exist_ok=True)
        
        for page_idx, page in enumerate(reader.pages):
            images = page.images
            for img_idx, img in enumerate(images):
                filename = f"page_{page_idx + 1}_img_{img_idx + 1}.{img.name.split('.')[-1] if '.' in img.name else 'png'}"
                filepath = os.path.join(output_dir, filename)
                
                with open(filepath, "wb") as f:
                    f.write(img.data)
                
                image_info.append({
                    "page": page_idx + 1,
                    "image_index": img_idx + 1,
                    "filename": filename,
                    "format": filename.split('.')[-1],
                    "width": img.image.width if hasattr(img, 'image') else 0,
                    "height": img.image.height if hasattr(img, 'image') else 0
                })
        return image_info

    @staticmethod
    def execute_page_operations(pdf_path: str, operations: List[Dict[str, Any]], output_path: str):
        """
        Execute page operations: rotate, delete, duplicate, reorder using pypdf.
        operations format: [
            {"type": "rotate", "page": 1, "angle": 90},
            {"type": "delete", "page": 2},
            {"type": "duplicate", "page": 3},
            {"type": "reorder", "order": [1, 3, 2]} # 1-indexed page list
        ]
        """
        reader = pypdf.PdfReader(pdf_path)
        writer = pypdf.PdfWriter()
        
        # 1. Determine the page sequence
        page_sequence = list(range(len(reader.pages)))
        
        # Check if there is a reorder operation first
        reorder_op = next((op for op in operations if op.get("type") == "reorder"), None)
        if reorder_op:
            page_sequence = [p - 1 for p in reorder_op["order"] if 0 <= p - 1 < len(reader.pages)]
            
        # Apply duplicates and deletes in sequence
        target_sequence = list(page_sequence)
        
        # We need to process structural operations in a stable way
        # Since deletes change index, we can build the list by mapping operations
        for op in operations:
            op_type = op.get("type")
            if op_type == "delete":
                p_idx = op["page"] - 1
                if p_idx in target_sequence:
                    # Remove the first occurrence
                    target_sequence.remove(p_idx)
            elif op_type == "duplicate":
                p_idx = op["page"] - 1
                if p_idx in target_sequence:
                    pos = target_sequence.index(p_idx)
                    target_sequence.insert(pos + 1, p_idx)
                    
        # Add pages to the writer
        for p_idx in target_sequence:
            writer.add_page(reader.pages[p_idx])
            
        # 2. Apply rotations on the writer pages
        for op in operations:
            op_type = op.get("type")
            if op_type == "rotate":
                # Find page in the new target sequence to rotate it
                p_idx = op["page"] - 1
                if 0 <= p_idx < len(writer.pages):
                    page = writer.pages[p_idx]
                    current_rotation = page.get("/Rotate", 0)
                    target_rotation = (current_rotation + op["angle"]) % 360
                    page.rotate(target_rotation)
                    
        # Save output
        with open(output_path, "wb") as out_f:
            writer.write(out_f)

    @staticmethod
    def replace_text(pdf_path: str, replacements: List[Dict[str, Any]], output_path: str):
        """
        Text replacements are handled client-side using pdf-lib for layout fidelity.
        This backend placeholder copies the PDF.
        """
        shutil.copy2(pdf_path, output_path)

    @staticmethod
    def compress_pdf(pdf_path: str, output_path: str):
        """Compress PDF content streams using pypdf."""
        reader = pypdf.PdfReader(pdf_path)
        writer = pypdf.PdfWriter()
        
        for page in reader.pages:
            page.compress_content_streams()
            writer.add_page(page)
            
        writer.add_metadata(reader.metadata)
        with open(output_path, "wb") as f:
            writer.write(f)

    @staticmethod
    def add_watermark(pdf_path: str, watermark_text: str, output_path: str):
        """
        Watermarking is handled client-side using pdf-lib.
        Backend placeholder copies the PDF.
        """
        shutil.copy2(pdf_path, output_path)

    @staticmethod
    def export_to_docx(pdf_path: str, output_path: str):
        """Convert PDF text to a Word Document (DOCX)."""
        reader = pypdf.PdfReader(pdf_path)
        word_doc = Document()
        word_doc.add_heading("Exported PDF Document", level=0)
        
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            paragraphs = text.split("\n\n")
            for p_text in paragraphs:
                p_text_clean = p_text.strip().replace("\n", " ")
                if p_text_clean:
                    word_doc.add_paragraph(p_text_clean)
            if i < len(reader.pages) - 1:
                word_doc.add_page_break()
            
        word_doc.save(output_path)

    @staticmethod
    def export_to_pptx(pdf_path: str, output_path: str):
        """Convert PDF to PowerPoint presentation."""
        reader = pypdf.PdfReader(pdf_path)
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        blank_slide_layout = prs.slide_layouts[6] # Blank
        
        for page in reader.pages:
            slide = prs.slides.add_slide(blank_slide_layout)
            text = (page.extract_text() or "").strip()
            
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(4.625))
            tf = txBox.text_frame
            tf.word_wrap = True
            
            lines = text.split("\n")
            if lines:
                p = tf.paragraphs[0]
                p.text = lines[0]
                for line in lines[1:]:
                    p = tf.add_paragraph()
                    p.text = line
                    p.font.size = Pt(14)
                
        prs.save(output_path)

    @staticmethod
    def export_to_excel(pdf_path: str, output_path: str):
        """Extract tables and save into Excel workbook using openpyxl directly (no pandas/numpy required)."""
        tables = PDFProcessor.extract_tables(pdf_path)
        
        import openpyxl
        wb = openpyxl.Workbook()
        
        if not tables:
            ws = wb.active
            ws.title = "Sheet1"
            ws.append(["No tables found in PDF"])
        else:
            # Remove default active sheet and create custom sheets per table
            default_sheet = wb.active
            wb.remove(default_sheet)
            
            for t in tables:
                sheet_name = f"Page_{t['page']}_Table_{t['table_index']}"[:30]
                ws = wb.create_sheet(title=sheet_name)
                for row in t["data"]:
                    ws.append(row)
                    
        wb.save(output_path)
